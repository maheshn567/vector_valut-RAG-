from fastapi import APIRouter, HTTPException, Request, UploadFile,Depends
from typing import Optional, List
from app.models.requests import ExtractionRequest
from app.models.response import ExtractionResponse, ExtractedPage
import pymupdf as fitz 
import docx 
import pptx 
from bs4 import BeautifulSoup
import requests
import json
import trafilatura
from app.utility.security import validate_api_key

router = APIRouter()

@router.post("/extract")
async def extract_content(request: Request,api_key: str = Depends(validate_api_key)):
    content_type = request.headers.get("content-type", "")
    
    file_bytes = None
    filename = None
    file_content_type = None
    url = None
    text = None
    webpage_title = None
    
    if "multipart/form-data" in content_type:
        form = await request.form()
        form_file = form.get("file")
        if form_file and isinstance(form_file, UploadFile):
            file_bytes = await form_file.read()
            filename = form_file.filename
            file_content_type = form_file.content_type
        
        url = form.get("url")
        text = form.get("text")
        
    elif "application/json" in content_type:
        try:
            body = await request.json()
            url = body.get("url")
            text = body.get("text")
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON body")
    else:
        # Fallback to query params
        url = request.query_params.get("url")
        text = request.query_params.get("text")

    # Handle the extracted inputs
    pages: List[ExtractedPage] = []
    doc_name = "unknown"
    doc_type = "unknown"
    
    if file_bytes is not None:
        doc_name = filename or "uploaded_file"
        doc_type = file_content_type or "application/octet-stream"
        
        # Parse PDF using PyMuPDF (fitz)
        if doc_type == "application/pdf" or doc_name.endswith(".pdf"):
            try:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                for i, page in enumerate(doc):
                    pages.append(ExtractedPage(page_number=i + 1, text=page.get_text()))
                doc_type = "pdf"
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Failed to parse PDF: {str(e)}")
                
        # Parse Word Doc
        elif doc_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"] or doc_name.endswith(".docx"):
            try:
                from io import BytesIO
                doc = docx.Document(BytesIO(file_bytes))
                full_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                pages.append(ExtractedPage(page_number=1, text=full_text))
                doc_type = "docx"
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Failed to parse DOCX: {str(e)}")
                
        # Parse PowerPoint
        elif doc_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"] or doc_name.endswith(".pptx"):
            try:
                from io import BytesIO
                prs = pptx.Presentation(BytesIO(file_bytes))
                slide_texts = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    slide_texts.append("\n".join(slide_text))
                for i, stext in enumerate(slide_texts):
                    pages.append(ExtractedPage(page_number=i + 1, text=stext))
                doc_type = "pptx"
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Failed to parse PPTX: {str(e)}")
                
        # Parse Plain Text
        elif doc_type == "text/plain" or doc_name.endswith(".txt"):
            try:
                text_content = file_bytes.decode("utf-8")
                pages.append(ExtractedPage(page_number=1, text=text_content))
                doc_type = "text"
            except UnicodeDecodeError:
                # Fallback decoding
                text_content = file_bytes.decode("latin1")
                pages.append(ExtractedPage(page_number=1, text=text_content))
                doc_type = "text"
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format")

    elif url:
        doc_name = url
        doc_type = "html"
        try:
            # 1. Fetch webpage raw HTML using trafilatura
            downloaded = trafilatura.fetch_url(url)
            if downloaded is None:
                # Fallback download using requests if trafilatura fails to fetch
                headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
                resp = requests.get(url, headers=headers, timeout=10)
                resp.raise_for_status()
                downloaded = resp.text
                
            # 2. Extract clean content and metadata
            extracted = trafilatura.bare_extraction(downloaded, as_dict=True)
            if extracted is None or not extracted.get("text"):
                # Fallback to BeautifulSoup if trafilatura extraction fails
                soup = BeautifulSoup(downloaded, "html.parser")
                for script in soup(["script", "style"]):
                    script.decompose()
                text_content = soup.get_text(separator="\n")
                cleaned_lines = [line.strip() for line in text_content.splitlines() if line.strip()]
                text_content = "\n".join(cleaned_lines)
                webpage_title = doc_name.split("/")[-1] if "/" in doc_name else doc_name
            else:
                text_content = extracted["text"]
                webpage_title = extracted.get("title") or (doc_name.split("/")[-1] if "/" in doc_name else doc_name)

            pages.append(ExtractedPage(page_number=1, text=text_content))
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to fetch or parse URL: {str(e)}")
            
    elif text:
        doc_name = "raw_text"
        doc_type = "text"
        pages.append(ExtractedPage(page_number=1, text=text))
        
    else:
        return {
            "success": False,
            "data": None,
            "error": {
                "code": "INVALID_INPUT",
                "message": "You must provide either a 'file', a 'url', or 'text'."
            },
            "metadata": {}
        }
        
    extracted_doc = ExtractionResponse(
        document_name=doc_name,
        content_type=doc_type,
        pages=pages,
        metadata={
            "page_count": len(pages),
            "title": webpage_title if url else (doc_name.split("/")[-1] if "/" in doc_name else doc_name)
        }
    )
    
    return {
        "success": True,
        "data": extracted_doc.model_dump(),
        "error": None,
        "metadata": {}
    }